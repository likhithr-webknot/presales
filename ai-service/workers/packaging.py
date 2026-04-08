"""
Packaging Agent
Assembles research + context outputs into a final PPTX or DOCX deliverable.
Uploads to MinIO and returns a presigned URL.

Stage 1: 5-7 slide PPTX (First Meeting Deck)
Stage 2: 8-10 slide PPTX (Post-Discovery Deck)
"""
import io
import json
import logging
import os
import re
from typing import Any, Optional
from datetime import datetime, timedelta

from minio import Minio
from minio.error import S3Error
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from config import get_settings
from schemas.agents import PackagingInput, PackagingOutput, ResearchBrief, WebknotContextOutput

logger = logging.getLogger(__name__)

# ── Webknot brand colours ─────────────────────────────────────────────────────
BRAND_DARK_NAVY  = RGBColor(0x0A, 0x1A, 0x3D)   # #0A1A3D — header backgrounds
BRAND_CYAN       = RGBColor(0x00, 0xC9, 0xFF)   # #00C9FF — accent / highlights
BRAND_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)

QUALITY_BANNED = [
    r"\{\{", r"\[\[", r"\bTBD\b", r"\bINSERT\b", r"\bPLACEHOLDER\b",
    r"\bLOREM\b", r"\bipsum\b", r"<[^>]+>",
]

SLIDE_PLAN_PROMPT = """You are a presales deck designer for Webknot Technologies.
Given research about a prospect and Webknot's positioning, create a slide plan for a first meeting deck.

Return ONLY valid JSON (no markdown fences) matching this schema:
{
  "slides": [
    {
      "title": "slide title",
      "type": "title|agenda|problem|solution|capability|case_study|next_steps",
      "bullets": ["bullet point 1", "bullet point 2", "bullet point 3"],
      "speaker_notes": "brief speaker guidance"
    }
  ]
}

Rules:
- 5-7 slides for Stage 1 (first meeting), 8-10 for Stage 2 (post-discovery)
- First slide = title slide (company name + tagline + date)
- Last slide = clear next steps
- Bullets: max 4 per slide, each under 12 words
- Be specific to the prospect — no generic filler
- Speaker notes: 1-2 sentences of guidance
"""


def _quality_check(text: str) -> list[str]:
    """Scan text for placeholder patterns. Returns list of warnings."""
    warnings = []
    for pattern in QUALITY_BANNED:
        if re.search(pattern, text, re.IGNORECASE):
            warnings.append(f"Placeholder text detected matching pattern: {pattern}")
    return warnings


def _add_title_slide(prs: Presentation, client_name: str, date_str: str) -> None:
    """Add branded title slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_DARK_NAVY

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Webknot Technologies × {client_name}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRAND_WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Partnership Overview | {date_str}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = BRAND_CYAN


def _add_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str = "") -> None:
    """Add a standard content slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    hfill = header.fill
    hfill.solid()
    hfill.fore_color.rgb = BRAND_DARK_NAVY
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BRAND_WHITE

    # Bullet content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    ctf = content_box.text_frame
    ctf.word_wrap = True

    for i, bullet in enumerate(bullets[:5]):
        if i == 0:
            para = ctf.paragraphs[0]
        else:
            para = ctf.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(16)
        para.space_after = Pt(8)

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


async def _generate_slide_plan(
    inp: PackagingInput,
    research: Optional[ResearchBrief],
    context: Optional[WebknotContextOutput],
    settings,
) -> list[dict]:
    """Ask LLM to create a slide-by-slide plan before we build the PPTX."""
    openai_client = AsyncOpenAI(api_key=settings.openai_api_key)

    slide_count = "5-7" if inp.stage == 1 else "8-10"

    research_summary = ""
    if research:
        research_summary = (
            f"Company: {research.company.name}\n"
            f"Challenges: {', '.join(research.company.key_challenges[:4])}\n"
            f"Talking points: {', '.join(research.talking_points[:4])}\n"
            f"Industry trends: {', '.join(research.industry.trends[:3])}"
        )

    context_summary = ""
    if context:
        context_summary = (
            f"Positioning: {context.positioning_narrative[:300]}\n"
            f"Differentiators: {', '.join(context.differentiators[:4])}"
        )

    user_prompt = (
        f"Create a {slide_count} slide deck for Stage {inp.stage} ({inp.collateral_type}).\n"
        f"Client: {inp.additional_context.get('client_name', 'the client')}\n"
        f"Domain: {inp.additional_context.get('domain', 'technology')}\n\n"
        f"Research Summary:\n{research_summary or 'Not available'}\n\n"
        f"Webknot Positioning:\n{context_summary or 'Not available'}"
    )

    response = await openai_client.chat.completions.create(
        model=settings.llm_cheap_model,
        messages=[
            {"role": "system", "content": SLIDE_PLAN_PROMPT},
            {"role": "user", "content": user_prompt},
        ],
        max_tokens=1200,
        temperature=0.3,
        response_format={"type": "json_object"},
    )

    raw = response.choices[0].message.content or '{"slides":[]}'
    try:
        data = json.loads(raw)
        return data.get("slides", [])
    except json.JSONDecodeError:
        logger.error("Slide plan LLM returned invalid JSON: %s", raw[:200])
        return []


def _get_minio_client() -> Minio:
    """Create a MinIO client from environment variables."""
    endpoint = os.getenv("STORAGE_ENDPOINT", "minio:9000").replace("http://", "").replace("https://", "")
    access_key = os.getenv("STORAGE_ACCESS_KEY", "minioadmin")
    secret_key = os.getenv("STORAGE_SECRET_KEY", "minioadmin")
    secure = os.getenv("STORAGE_ENDPOINT", "").startswith("https")
    return Minio(endpoint, access_key=access_key, secret_key=secret_key, secure=secure)


def _upload_to_minio_sync(file_bytes: bytes, key: str, content_type: str) -> str:
    """
    Upload file to MinIO bucket and return a presigned URL (24h TTL).
    Synchronous — called from async context via run_in_executor if needed.
    """
    bucket = os.getenv("STORAGE_BUCKET_ARTIFACTS", "presales-artifacts")
    client = _get_minio_client()

    # Ensure bucket exists
    try:
        if not client.bucket_exists(bucket):
            client.make_bucket(bucket)
    except S3Error as exc:
        logger.warning("MinIO bucket check/create failed: %s", exc)

    # Upload
    try:
        client.put_object(
            bucket,
            key,
            io.BytesIO(file_bytes),
            length=len(file_bytes),
            content_type=content_type,
        )
        # Presigned URL — 24 hour TTL
        url = client.presigned_get_object(bucket, key, expires=timedelta(hours=24))
        logger.info("MinIO upload success: bucket=%s key=%s size=%d", bucket, key, len(file_bytes))
        return url
    except S3Error as exc:
        logger.error("MinIO upload failed for key=%s: %s", key, exc)
        # Return a fallback key-based URL so the job doesn't fully fail
        return f"/api/artifacts/{key}"


async def _upload_to_minio(file_bytes: bytes, key: str, content_type: str) -> str:
    """Async wrapper around synchronous MinIO upload."""
    import asyncio
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _upload_to_minio_sync, file_bytes, key, content_type)


async def run(payload: dict[str, Any], engagement_id: Optional[str]) -> dict[str, Any]:
    """
    Packaging Agent entry point.
    Assembles research + context into a PPTX and uploads to MinIO.
    """
    inp = PackagingInput(
        engagement_id=engagement_id or payload.get("engagement_id", ""),
        collateral_type=payload.get("collateral_type", "FIRST_MEETING_DECK"),
        stage=payload.get("stage", 1),
        research_brief=ResearchBrief(**payload["research_brief"]) if payload.get("research_brief") else None,
        webknot_context=WebknotContextOutput(**payload["webknot_context"]) if payload.get("webknot_context") else None,
        additional_context=payload.get("additional_context", {}),
        output_format=payload.get("output_format", "pptx"),
        version=payload.get("version", 1),
    )

    settings = get_settings()
    client_name = inp.additional_context.get("client_name", "Client")
    date_str = datetime.utcnow().strftime("%B %Y")

    logger.info(
        "Packaging agent starting: engagement=%s type=%s stage=%d format=%s",
        inp.engagement_id, inp.collateral_type, inp.stage, inp.output_format,
    )

    # ── Generate slide plan ────────────────────────────────────────────────────
    slides = await _generate_slide_plan(inp, inp.research_brief, inp.webknot_context, settings)

    if not slides:
        logger.warning("Slide plan generation returned empty — using minimal fallback")
        slides = [
            {"title": f"Webknot × {client_name}", "type": "title", "bullets": [], "speaker_notes": ""},
            {"title": "About Webknot", "type": "capability", "bullets": ["Custom software development", "AI/ML integration", "Cloud & DevOps", "Data engineering"], "speaker_notes": ""},
            {"title": "Why We're Here", "type": "problem", "bullets": ["Understanding your challenges", "Exploring partnership opportunities"], "speaker_notes": ""},
            {"title": "Next Steps", "type": "next_steps", "bullets": ["Follow-up call", "Requirements deep dive", "Proposal preparation"], "speaker_notes": ""},
        ]

    # ── Build PPTX ────────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    quality_warnings: list[str] = []

    for i, slide_data in enumerate(slides):
        title = slide_data.get("title", f"Slide {i+1}")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("speaker_notes", "")

        # Quality check all text
        all_text = " ".join([title] + bullets + [notes])
        quality_warnings.extend(_quality_check(all_text))

        if i == 0 and slide_data.get("type") == "title":
            _add_title_slide(prs, client_name, date_str)
        else:
            _add_content_slide(prs, title, bullets, notes)

    # ── Serialize to bytes ─────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    # ── Upload to MinIO ────────────────────────────────────────────────────────
    file_key = f"{inp.engagement_id}/v{inp.version}/{inp.collateral_type.lower()}.pptx"
    presigned_url = await _upload_to_minio(pptx_bytes, file_key, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    output = PackagingOutput(
        engagement_id=inp.engagement_id,
        file_key=file_key,
        presigned_url=presigned_url,
        format="pptx",
        slide_count=len(prs.slides),
        quality_warnings=list(set(quality_warnings)),
        version=inp.version,
    )

    logger.info(
        "Packaging complete: engagement=%s slides=%d quality_warnings=%d",
        inp.engagement_id, len(prs.slides), len(output.quality_warnings),
    )

    return output.model_dump()
